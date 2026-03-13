from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Colors
BG = RGBColor(0x0f, 0x0f, 0x0f)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xAA, 0xAA, 0xAA)
ACCENT = RGBColor(0x7C, 0xE8, 0xA0)   # green accent
CODE_BG = RGBColor(0x1E, 0x1E, 0x1E)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_textbox(slide, text, left, top, width, height,
                font_size=28, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, font_name="Helvetica Neue",
                wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_code_block(slide, code, left, top, width, height, font_size=14):
    # dark box
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RGBColor(0x33, 0x33, 0x33)

    txBox = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), height - Inches(0.3)
    )
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code.strip()
    run.font.size = Pt(font_size)
    run.font.color.rgb = ACCENT
    run.font.name = "Courier New"


def title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, title,
                Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
                font_size=52, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, subtitle,
                Inches(1), Inches(4.2), Inches(11.3), Inches(0.8),
                font_size=24, color=GRAY, align=PP_ALIGN.CENTER)
    return slide


def content_slide(prs, title, body=None, code=None, code_size=14):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # title
    add_textbox(slide, title,
                Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9),
                font_size=36, bold=True)
    # divider line
    from pptx.util import Pt as UPt
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.35),
                                  Inches(11.9), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    y_offset = Inches(1.5)

    if body:
        add_textbox(slide, body,
                    Inches(0.7), y_offset, Inches(11.9), Inches(1.0),
                    font_size=22, color=GRAY)
        y_offset += Inches(1.1)

    if code:
        code_h = min(Inches(5.2), Inches(0.5 + 0.28 * (code.count('\n') + 1)))
        add_code_block(slide, code,
                       Inches(0.7), y_offset, Inches(11.9), code_h,
                       font_size=code_size)
    return slide


def two_col_slide(prs, title, left_col, right_col, body=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, title,
                Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9),
                font_size=36, bold=True)
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.35),
                                  Inches(11.9), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    if body:
        add_textbox(slide, body,
                    Inches(0.7), Inches(1.5), Inches(11.9), Inches(0.8),
                    font_size=20, color=GRAY)

    # Left header
    add_textbox(slide, left_col[0],
                Inches(0.7), Inches(2.4), Inches(5.5), Inches(0.6),
                font_size=20, bold=True, color=ACCENT)
    for i, line_text in enumerate(left_col[1:]):
        add_textbox(slide, line_text,
                    Inches(0.7), Inches(3.0 + i * 0.6), Inches(5.5), Inches(0.6),
                    font_size=18, color=GRAY)

    # Right header
    add_textbox(slide, right_col[0],
                Inches(7.0), Inches(2.4), Inches(5.5), Inches(0.6),
                font_size=20, bold=True, color=ACCENT)
    for i, line_text in enumerate(right_col[1:]):
        add_textbox(slide, line_text,
                    Inches(7.0), Inches(3.0 + i * 0.6), Inches(5.5), Inches(0.6),
                    font_size=18, color=GRAY)
    return slide


def three_prop_slide(prs, title, props):
    """props = [(name, desc), ...]"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_textbox(slide, title,
                Inches(0.7), Inches(0.4), Inches(11.9), Inches(0.9),
                font_size=36, bold=True)
    line = slide.shapes.add_shape(1, Inches(0.7), Inches(1.35),
                                  Inches(11.9), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    x_positions = [Inches(0.7), Inches(4.7), Inches(8.7)]
    for i, (name, desc) in enumerate(props):
        add_textbox(slide, name,
                    x_positions[i], Inches(2.8), Inches(3.6), Inches(0.9),
                    font_size=30, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        add_textbox(slide, desc,
                    x_positions[i], Inches(3.8), Inches(3.6), Inches(1.2),
                    font_size=17, color=GRAY, align=PP_ALIGN.CENTER)
    return slide


# ── Slides ────────────────────────────────────────────────────────────────────

# 1 — Title
title_slide(prs, "Agentic AI: From Chatbot to Agent", "A 50-minute workshop")

# 2 — What did you just see?
content_slide(prs, "What did you just see?",
              body="[show terminal output from opening demo]")

# 3 — The Chatbot Model
content_slide(prs, "The Chatbot Model",
              body="You are the agent. The LLM reacts.",
              code=(
                  "[ You ] ──── question ────► [ LLM ]\n"
                  "[ You ] ◄─── answer ──────  [ LLM ]\n"
                  "[ You ] ──── question ────► [ LLM ]\n"
                  "[ You ] ◄─── answer ──────  [ LLM ]"
              ), code_size=20)

# 4 — The Agent Loop
content_slide(prs, "The Agent Loop",
              code=(
                  "       PERCEIVE\n"
                  "    (task + context)\n"
                  "           │\n"
                  "           ▼\n"
                  "        REASON          ◄───┐\n"
                  "     (LLM decides)          │\n"
                  "           │                │\n"
                  "           ▼                │\n"
                  "          ACT           OBSERVE\n"
                  "     (tool call)    (tool result back)\n"
                  "           │                │\n"
                  "           └────────────────┘"
              ), code_size=16)

# 5 — The Intern Analogy
two_col_slide(prs, "The LLM is a reasoning step, not a runtime",
              left_col=["Smart Intern",
                        "Reads the task",
                        'Writes "please search X"',
                        "Reads the result",
                        'Writes "save this summary"',
                        'Writes "done"'],
              right_col=["Your Code",
                         "Sends the task",
                         "Runs the search",
                         "Returns the result",
                         "Saves the file",
                         "Stops the loop"],
              body="The intern never touched the keyboard.")

# 6 — What Makes It Agentic
three_prop_slide(prs, "Three properties of an agent",
                 [
                     ("AUTONOMY", "completes sub-tasks without step-by-step direction"),
                     ("TOOL USE", "can reach outside the model — search, write, call APIs"),
                     ("STATE", "carries context across multiple reasoning steps"),
                 ])

# 7 — Why Now?
content_slide(prs, "Why is this possible now?",
              body="Before (2022): parse free text and hope.\nAfter (2023+): structured JSON — reliable, parseable, automatable.",
              code='{ "tool": "search_web", "query": "agentic AI" }',
              code_size=20)

# 8 — Demo 1 Setup
content_slide(prs, "Demo 1: The Reasoning Illusion",
              body="Can Claude complete a multi-step task... without tools?")

# 9 — Demo 1 Debrief
content_slide(prs, "Reasoning ≠ Action",
              body="Very smart autocomplete. Nothing more.",
              code=(
                  "[✓] Reasoned about the task\n"
                  "[✓] Generated a plan\n"
                  "[✓] Wrote 'search results' (hallucinated)\n"
                  "[✗] Searched anything\n"
                  "[✗] Saved any file\n"
                  "[✗] Completed the task"
              ), code_size=18)

# 10 — Demo 2 Setup
content_slide(prs, "Demo 2: Giving the Model Hands",
              body="What happens when Claude knows it has tools?",
              code=(
                  "You have access to these tools:\n"
                  "- search_web(query: str)\n"
                  "- save_note(text: str, filename: str)\n\n"
                  "When you want to use a tool, respond with:\n"
                  "TOOL_CALL: tool_name(arg1, arg2)"
              ), code_size=16)

# 11 — Demo 2 Debrief
content_slide(prs, "Tool use is structured conversation",
              body='You just played the role of agent.py.',
              code=(
                  'You:     "Research agentic AI and save a summary."\n'
                  'Claude:  TOOL_CALL: search_web("agentic AI definition")\n'
                  'You:     [result: "Agentic AI refers to..."]\n'
                  'Claude:  TOOL_CALL: save_note("• Point 1...", "summary.txt")\n'
                  'You:     [result: "Saved to notes/summary.txt"]\n'
                  'Claude:  "Done. Summary saved."'
              ), code_size=15)

# 12 — Demo 3 Setup
content_slide(prs, "Demo 3: The Full Loop",
              body="Now the code plays the role you just played.",
              code="python agent.py", code_size=28)

# 13 — Code: TOOLS Schema
content_slide(prs, "The Tool Menu",
              body="The description is read by the model, not the computer.",
              code=(
                  'TOOLS = [{\n'
                  '    "name": "search_web",\n'
                  '    "description": "Search DuckDuckGo and return a summary.",\n'
                  '    "input_schema": {\n'
                  '        "properties": {\n'
                  '            "query": {"type": "string"}\n'
                  '        }\n'
                  '    }\n'
                  '}]'
              ), code_size=15)

# 14 — Code: The Loop
content_slide(prs, "The Whole Agent — 15 Lines",
              body="Everything else is plumbing.",
              code=(
                  "while iteration < MAX_ITERATIONS:\n"
                  "    response = client.messages.create(\n"
                  "        model=MODEL, tools=TOOLS, messages=messages\n"
                  "    )\n"
                  "    if response.stop_reason == 'end_turn':\n"
                  "        break\n"
                  "    if response.stop_reason == 'tool_use':\n"
                  "        result = execute_tool(tool_call)\n"
                  "        messages.append(tool_result)"
              ), code_size=15)

# 15 — Multi-Agent Systems
content_slide(prs, "One agent is just the start",
              body="Specialists. Parallelism. One agent checks another.",
              code=(
                  "         [Orchestrator]\n"
                  "        /       |       \\\n"
                  "[Research]  [Critic]  [Writer]\n"
                  "   agent      agent     agent"
              ), code_size=20)

# 16 — Where It's Going
content_slide(prs, "Three things to watch",
              body="The loop is the same. The tools get more powerful.",
              code=(
                  "⏳  Long-horizon tasks — agents that run for hours\n"
                  "🖥️  Computer use — agents that control a browser or desktop\n"
                  "🤝  Agent-to-agent protocols — agents from different companies collaborating"
              ), code_size=16)

# 17 — Clone the Repo
content_slide(prs, "Your turn",
              body="Running in under 10 minutes.",
              code=(
                  "git clone [repo URL]\n"
                  "cd repo\n"
                  "pip install anthropic python-dotenv requests\n"
                  "cp .env.example .env\n"
                  "# Add your API key\n"
                  "python agent.py"
              ), code_size=18)

# 18 — Q&A
content_slide(prs, "Questions?",
              body=(
                  "• What task would you want to automate with an agent?\n"
                  "• What could go wrong if an agent had too many powerful tools?\n"
                  "• How would you test an agent?"
              ))

out = "/Users/bhagatpranish/Documents/daisy-workshop/workshop/slides.pptx"
prs.save(out)
print(f"Saved → {out}")
